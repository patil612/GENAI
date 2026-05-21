"""
Universal File Parser Service
Supports: PDF, DOCX, DOC, TXT, RTF, ODT, XLSX, XLS, CSV, PPTX, PPT, HTML, MD, images (OCR)
"""

import os
import io

def extract_text_from_file(file_path: str, filename: str) -> str:
    """
    Dispatch to the correct parser based on file extension.
    Returns extracted plain text, or raises an exception on failure.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".docx",):
        return _parse_docx(file_path)
    elif ext in (".doc",):
        return _parse_doc_mammoth(file_path)
    elif ext in (".txt", ".md", ".log", ".ini", ".cfg", ".json", ".xml", ".yaml", ".yml"):
        return _parse_plain_text(file_path)
    elif ext == ".rtf":
        return _parse_rtf(file_path)
    elif ext == ".odt":
        return _parse_odt(file_path)
    elif ext in (".xlsx", ".xls"):
        return _parse_excel(file_path, ext)
    elif ext == ".csv":
        return _parse_csv(file_path)
    elif ext in (".pptx",):
        return _parse_pptx(file_path)
    elif ext in (".ppt",):
        return _parse_ppt_fallback(file_path)
    elif ext in (".html", ".htm"):
        return _parse_html(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"):
        return _parse_image_ocr(file_path)
    else:
        # Generic fallback — try raw UTF-8 read
        return _parse_plain_text(file_path)


# ─── Individual Parsers ────────────────────────────────────────────────────────

def _parse_pdf(path: str) -> str:
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n".join(p for p in pages if p.strip())
        if not text.strip():
            # Fallback: try pdfplumber for scanned/complex PDFs
            try:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    text = "\n".join(
                        page.extract_text() or "" for page in pdf.pages
                    )
            except ImportError:
                pass
        return text or ""
    except Exception as e:
        raise ValueError(f"PDF parsing failed: {e}")


def _parse_docx(path: str) -> str:
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"DOCX parsing failed: {e}")


def _parse_doc_mammoth(path: str) -> str:
    """Legacy .doc format via mammoth (converts to plain text)."""
    try:
        import mammoth
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
        return result.value or ""
    except ImportError:
        # If mammoth not installed, try plain read
        return _parse_plain_text(path)
    except Exception as e:
        raise ValueError(f"DOC parsing failed: {e}")


def _parse_plain_text(path: str) -> str:
    encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc, errors="strict") as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    # Last resort: ignore errors
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_rtf(path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return rtf_to_text(content)
    except ImportError:
        # Fallback: crude strip of RTF tags
        import re
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        text = re.sub(r'\{[^{}]*\}', '', raw)
        text = re.sub(r'\\[a-zA-Z]+\d*\s?', '', text)
        return text.strip()
    except Exception as e:
        raise ValueError(f"RTF parsing failed: {e}")


def _parse_odt(path: str) -> str:
    try:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype
        doc = load(path)
        paragraphs = doc.getElementsByType(P)
        return "\n".join(teletype.extractText(p) for p in paragraphs)
    except ImportError:
        raise ValueError("ODT support requires 'odfpy'. Install via: pip install odfpy")
    except Exception as e:
        raise ValueError(f"ODT parsing failed: {e}")


def _parse_excel(path: str, ext: str) -> str:
    try:
        import openpyxl
        if ext == ".xlsx":
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append(" | ".join(cells))
            return "\n".join(rows)
        else:
            # .xls — use xlrd
            import xlrd
            wb = xlrd.open_workbook(path)
            rows = []
            for sheet in wb.sheets():
                rows.append(f"[Sheet: {sheet.name}]")
                for rx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    rows.append(" | ".join(cells))
            return "\n".join(rows)
    except ImportError as e:
        raise ValueError(f"Excel support requires openpyxl/xlrd: {e}")
    except Exception as e:
        raise ValueError(f"Excel parsing failed: {e}")


def _parse_csv(path: str) -> str:
    import csv
    rows = []
    for enc in ["utf-8", "latin-1"]:
        try:
            with open(path, newline="", encoding=enc, errors="strict") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(row))
            return "\n".join(rows)
        except (UnicodeDecodeError, UnicodeError):
            rows = []
            continue
    return "\n".join(rows)


def _parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(slides_text)
    except ImportError:
        raise ValueError("PPTX support requires 'python-pptx'. Install via: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"PPTX parsing failed: {e}")


def _parse_ppt_fallback(path: str) -> str:
    """Legacy .ppt — try mammoth, else plain read."""
    try:
        import mammoth
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
        return result.value or ""
    except Exception:
        return _parse_plain_text(path)


def _parse_html(path: str) -> str:
    try:
        from html.parser import HTMLParser

        class _Stripper(HTMLParser):
            def __init__(self):
                super().__init__()
                self.fed = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    self.fed.append(data)
            def get_data(self):
                return " ".join(self.fed)

        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
        s = _Stripper()
        s.feed(html)
        return s.get_data()
    except Exception as e:
        raise ValueError(f"HTML parsing failed: {e}")


def _parse_image_ocr(path: str) -> str:
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except ImportError:
        raise ValueError(
            "Image OCR requires 'Pillow' and 'pytesseract'. "
            "Install via: pip install Pillow pytesseract  (and Tesseract OCR binary)"
        )
    except Exception as e:
        raise ValueError(f"OCR failed: {e}")
